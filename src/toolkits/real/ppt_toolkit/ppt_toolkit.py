from langchain_core.tools import tool
from mcp.server.fastmcp import FastMCP
import pptx
import os
from typing import Optional


mcp = FastMCP('ppt_toolkit')
def ppt_create_presentation(file_path: str) -> dict:
    """
    Create a new PowerPoint presentation.

    Args:
        file_path (str): The path to save the presentation.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    dir_name = os.path.dirname(file_path)
    if not os.path.exists(dir_name):
        os.makedirs(dir_name, exist_ok=True)
    prs = pptx.Presentation()
    prs.save(file_path)
    return {'success': True, 'message': f'Presentation created at {file_path}'}


@tool(parse_docstring=True)
def langgraph_ppt_create_presentation(file_path: str) -> dict:
    """
    Create a new PowerPoint presentation.

    Args:
        file_path (str): The path to save the presentation.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_create_presentation(file_path=file_path)


@mcp.tool()
def mcp_ppt_create_presentation(file_path: str) -> dict:
    """
    Create a new PowerPoint presentation.

    Args:
        file_path (str): The path to save the presentation.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_create_presentation(file_path=file_path)


def ppt_delete_slide(file_path: str, slide_index: int) -> dict:
    """
    Delete a slide from the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to delete.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    if not os.path.exists(file_path):
        return {'success': False, 'message': f'File {file_path} does not exist.'}
    prs = pptx.Presentation(file_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {'success': False, 'message': f'Slide index {slide_index} is out of range.'}
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[slide_index])
    prs.save(file_path)
    return {'success': True, 'message': f'Slide {slide_index} deleted from {file_path}'}


@tool(parse_docstring=True)
def langgraph_ppt_delete_slide(file_path: str, slide_index: int) -> dict:
    """
    Delete a slide from the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to delete.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_delete_slide(file_path=file_path, slide_index=slide_index)


@mcp.tool()
def mcp_ppt_delete_slide(file_path: str, slide_index: int) -> dict:
    """
    Delete a slide from the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to delete.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_delete_slide(file_path=file_path, slide_index=slide_index)


def ppt_add_slide(file_path: str, title: str, content: str) -> dict:
    """
    Add a text slide to the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        title (str): The title of the slide.
        content (str): The content of the slide.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    if not os.path.exists(file_path):
        return {'success': False, 'message': f'File {file_path} does not exist.'}
    prs = pptx.Presentation(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    prs.save(file_path)
    return {'success': True, 'message': f'Text slide added to {file_path}'}


@tool(parse_docstring=True)
def langgraph_ppt_add_slide(file_path: str, title: str, content: str) -> dict:
    """
    Add a text slide to the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        title (str): The title of the slide.
        content (str): The content of the slide.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_add_slide(file_path=file_path, title=title, content=content)


@mcp.tool()
def mcp_ppt_add_slide(file_path: str, title: str, content: str) -> dict:
    """
    Add a text slide to the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        title (str): The title of the slide.
        content (str): The content of the slide.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_add_slide(file_path=file_path, title=title, content=content)


def ppt_edit_slide(file_path: str, slide_index: int, title: Optional[str]=None, content: Optional[str]=None) -> dict:
    """
    Edit a slide in the PowerPoint presentation.
    
    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to edit.
        title (Optional[str]): The new title for the slide. If None, the title will not be changed.
        content (Optional[str]): The new content for the slide. If None, the content will not be changed.
    
    Returns:
        dict: A dictionary containing the success status and message.
    """
    if not os.path.exists(file_path):
        return {'success': False, 'message': f'File {file_path} does not exist.'}
    prs = pptx.Presentation(file_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {'success': False, 'message': f'Slide index {slide_index} is out of range.'}
    slide = prs.slides[slide_index]
    if title is not None:
        slide.shapes.title.text = title
    if content is not None:
        slide.placeholders[1].text = content
    prs.save(file_path)
    return {'success': True, 'message': f'Slide {slide_index} edited in {file_path}'}


@tool(parse_docstring=True)
def langgraph_ppt_edit_slide(file_path: str, slide_index: int, title: Optional[str]=None, content: Optional[str]=None) -> dict:
    """
    Edit a slide in the PowerPoint presentation.
    
    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to edit.
        title (Optional[str]): The new title for the slide. If None, the title will not be changed.
        content (Optional[str]): The new content for the slide. If None, the content will not be changed.
    
    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_edit_slide(file_path=file_path, slide_index=slide_index, title=title, content=content)


@mcp.tool()
def mcp_ppt_edit_slide(file_path: str, slide_index: int, title: Optional[str]=None, content: Optional[str]=None) -> dict:
    """
    Edit a slide in the PowerPoint presentation.
    
    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to edit.
        title (Optional[str]): The new title for the slide. If None, the title will not be changed.
        content (Optional[str]): The new content for the slide. If None, the content will not be changed.
    
    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_edit_slide(file_path=file_path, slide_index=slide_index, title=title, content=content)


def ppt_add_picture(file_path: str, slide_index: int, picture_path: str) -> dict:
    """
    Add a picture to a slide in the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to add the picture to.
        picture_path (str): The path to the picture file.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    if not os.path.exists(file_path):
        return {'success': False, 'message': f'File {file_path} does not exist.'}
    if not os.path.exists(picture_path):
        return {'success': False, 'message': f'Picture {picture_path} does not exist.'}
    prs = pptx.Presentation(file_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {'success': False, 'message': f'Slide index {slide_index} is out of range.'}
    slide = prs.slides[slide_index]
    left = top = pptx.util.Inches(1)
    slide.shapes.add_picture(picture_path, left, top)
    prs.save(file_path)
    return {'success': True, 'message': f'Picture added to slide {slide_index} in {file_path}'}


@tool(parse_docstring=True)
def langgraph_ppt_add_picture(file_path: str, slide_index: int, picture_path: str) -> dict:
    """
    Add a picture to a slide in the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to add the picture to.
        picture_path (str): The path to the picture file.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_add_picture(file_path=file_path, slide_index=slide_index, picture_path=picture_path)


@mcp.tool()
def mcp_ppt_add_picture(file_path: str, slide_index: int, picture_path: str) -> dict:
    """
    Add a picture to a slide in the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to add the picture to.
        picture_path (str): The path to the picture file.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_add_picture(file_path=file_path, slide_index=slide_index, picture_path=picture_path)


