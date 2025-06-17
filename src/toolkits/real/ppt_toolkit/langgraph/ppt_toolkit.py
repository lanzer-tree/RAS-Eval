from langchain_core.tools import tool
from src.toolkits.real.ppt_toolkit.function.ppt_toolkit import ppt_create_presentation
from src.toolkits.real.ppt_toolkit.function.ppt_toolkit import ppt_delete_slide
from src.toolkits.real.ppt_toolkit.function.ppt_toolkit import ppt_add_slide
from src.toolkits.real.ppt_toolkit.function.ppt_toolkit import ppt_edit_slide
from src.toolkits.real.ppt_toolkit.function.ppt_toolkit import ppt_add_picture
import pptx
import os
from typing import Optional




@tool(parse_docstring=True)
def _ppt_create_presentation(file_path: str) -> dict:
    """
    Create a new PowerPoint presentation.

    Args:
        file_path (str): The path to save the presentation.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_create_presentation(file_path=file_path)


@tool(parse_docstring=True)
def _ppt_delete_slide(file_path: str, slide_index: int) -> dict:
    """
    Delete a slide from the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to delete.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_delete_slide(file_path=file_path, slide_index=slide_index)


@tool(parse_docstring=True)
def _ppt_add_slide(file_path: str, title: str, content: str) -> dict:
    """
    Add a text slide to the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        title (str): The title of the slide.
        content (str): The content of the slide.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_add_slide(file_path=file_path, title=title, content=content)


@tool(parse_docstring=True)
def _ppt_edit_slide(file_path: str, slide_index: int, title: Optional[str]=None, content: Optional[str]=None) -> dict:
    """
    Edit a slide in the PowerPoint presentation.
    
    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to edit.
        title (Optional[str]): The new title for the slide. If None, the title will not be changed.
        content (Optional[str]): The new content for the slide. If None, the content will not be changed.
    
    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_edit_slide(file_path=file_path, slide_index=slide_index, title=title, content=content)


@tool(parse_docstring=True)
def _ppt_add_picture(file_path: str, slide_index: int, picture_path: str) -> dict:
    """
    Add a picture to a slide in the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation.
        slide_index (int): The index of the slide to add the picture to.
        picture_path (str): The path to the picture file.

    Returns:
        dict: A dictionary containing the success status and message.
    """
    return ppt_add_picture(file_path=file_path, slide_index=slide_index, picture_path=picture_path)


